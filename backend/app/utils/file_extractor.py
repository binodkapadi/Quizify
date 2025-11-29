import io
import os
import shutil
import platform
from typing import List

from fastapi import UploadFile
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image

# Try to import pytesseract, but handle gracefully if not available
try:
    import pytesseract
    OCR_AVAILABLE = True

    # Detect OS
    system = platform.system()

    if system == "Linux":
        # ================================
        #   IMPORTANT FOR RENDER.COM
        #   Auto-detect Tesseract installed via apt
        # ================================
        tesseract_path = shutil.which("tesseract")
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
            print(f"✔ Using Tesseract at: {tesseract_path}")
        else:
            print("❌ Tesseract not found on Linux system")

    elif system == "Windows":
        # Windows auto-detection
        username = os.getenv('USERNAME', '')
        possible_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
            rf'C:\Users\{username}\AppData\Local\Programs\Tesseract-OCR\tesseract.exe' if username else None,
        ]

        possible_paths = [p for p in possible_paths if p]

        tesseract_in_path = shutil.which("tesseract")

        if tesseract_in_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_in_path
        else:
            for path in possible_paths:
                if os.path.exists(path):
                    pytesseract.pytesseract.tesseract_cmd = path
                    print(f"✔ Found Tesseract at: {path}")
                    break
            else:
                print("⚠️ Tesseract not found on Windows system")

except ImportError:
    OCR_AVAILABLE = False
    pytesseract = None


ALLOWED_TEXT_EXTENSIONS = {
    ".pdf", ".doc", ".docx", ".ppt", ".pptx", ".txt",
}

ALLOWED_IMAGE_EXTENSIONS = {
    ".png", ".jpg", ".jpeg",
}


async def extract_text_from_files(files: List[UploadFile]) -> str:
    """
    Extract text from uploaded files.
    Supports: PDF, Word, PowerPoint, plain text, and images (PNG, JPG, JPEG) via OCR.
    """
    texts: List[str] = []

    for file in files:
        filename = file.filename or ""
        ext = os.path.splitext(filename)[1].lower()

        raw = await file.read()
        if not raw:
            continue

        # Handle document files
        if ext in ALLOWED_TEXT_EXTENSIONS:
            if ext == ".pdf":
                try:
                    buffer = io.BytesIO(raw)
                    reader = PdfReader(buffer)
                    for page in reader.pages:
                        page_text = page.extract_text() or ""
                        texts.append(page_text)
                except Exception as e:
                    print(f"Error reading PDF {filename}: {e}")
                    continue

            elif ext in (".doc", ".docx"):
                try:
                    buffer = io.BytesIO(raw)
                    doc = Document(buffer)
                    texts.append("\n".join(p.text for p in doc.paragraphs))
                except Exception as e:
                    print(f"Error reading Word doc {filename}: {e}")
                    continue

            elif ext in (".ppt", ".pptx"):
                try:
                    buffer = io.BytesIO(raw)
                    presentation = Presentation(buffer)
                    slide_text = []
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slide_text.append(shape.text)
                    texts.append("\n".join(slide_text))
                except Exception as e:
                    print(f"Error reading PowerPoint {filename}: {e}")
                    continue

            elif ext == ".txt":
                try:
                    texts.append(raw.decode("utf-8", errors="ignore"))
                except Exception as e:
                    print(f"Error reading text file {filename}: {e}")
                    continue

        # Handle image files with OCR
        elif ext in ALLOWED_IMAGE_EXTENSIONS:
            if not OCR_AVAILABLE:
                raise ValueError(
                    "OCR support not available. Install pytesseract and Tesseract OCR."
                )
            try:
                image = Image.open(io.BytesIO(raw))
                if image.mode != "RGB":
                    image = image.convert("RGB")

                extracted_text = pytesseract.image_to_string(image, lang="eng")
                if extracted_text.strip():
                    texts.append(extracted_text.strip())
                else:
                    print(f"Warning: No text extracted from image {filename}")

            except Exception as e:
                error_msg = str(e)
                if "tesseract" in error_msg.lower():
                    raise ValueError(
                        "Tesseract OCR not found. Install Tesseract OCR on your system."
                    )
                print(f"Error extracting text from image {filename}: {e}")
                raise

    combined = "\n\n".join(t for t in texts if t.strip())
    return combined